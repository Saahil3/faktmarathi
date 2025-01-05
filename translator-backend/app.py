from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import requests
from PyPDF2 import PdfReader
import easyocr
from pptx import Presentation
from io import BytesIO

# Initialize Flask app
app = Flask(__name__)
from flask_cors import CORS

CORS(app, resources={r"/*": {"origins": ["https://faktmarathi.vercel.app"], "methods": ["GET", "POST", "OPTIONS"]}})

# Hugging Face API details
TRANSLATE_URL = "https://sepioo-facebook-translation.hf.space/translate"
HEADERS = {"Content-Type": "application/json"}

# Function to translate text using Hugging Face API
def translate_text(input_text, from_language="en", to_language="mr"):
    payload = {
        "from_language": from_language,
        "to_language": to_language,
        "input_text": input_text
    }

    response = requests.post(TRANSLATE_URL, json=payload, headers=HEADERS)

    if response.status_code == 200:
        result = response.json()
        translated_text = result.get("translate", "")
        if not translated_text:
            raise Exception("No translation available in the response.")
        return translated_text
    else:
        raise Exception(f"API Error: {response.status_code} - {response.text}")

# Function to extract text from a PDF
def extract_text_from_pdf(input_pdf):
    reader = PdfReader(input_pdf)
    content = ""
    for page in reader.pages:
        content += page.extract_text()
    if not content.strip():
        raise ValueError("The PDF appears to be empty or contains non-extractable text.")
    return content

# Function to extract text from an image using EasyOCR
def extract_text_from_image(image_file):
    reader = easyocr.Reader(['en'])  # Specify the language (English here)
    image_path = 'temp_image.png'  # Save the uploaded image temporarily
    image_file.save(image_path)
    result = reader.readtext(image_path)

    extracted_text = ' '.join([item[1] for item in result])
    if not extracted_text.strip():
        raise ValueError("No text found in the image.")
    return extracted_text

# Function to extract and translate text from a PPT file
def translate_ppt(input_ppt, from_language="en", to_language="mr"):
    presentation = Presentation(input_ppt)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                original_text = shape.text
                if original_text.strip():  # Translate only if text exists
                    translated_text = translate_text(original_text, from_language, to_language)
                    shape.text = translated_text
    # Save the modified PPT to a BytesIO object
    output_ppt = BytesIO()
    presentation.save(output_ppt)
    output_ppt.seek(0)  # Rewind to the beginning of the file-like object
    return output_ppt

# Endpoint for text translation
@app.route('/translate-text', methods=['POST'])
def translate_text_endpoint():
    data = request.json
    text = data.get('text', '')
    from_language = data.get('from_language', 'en')
    to_language = data.get('to_language', 'mr')

    try:
        translated_text = translate_text(text, from_language, to_language)
        return jsonify({'translated_text': translated_text})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Endpoint for document translation
@app.route('/translate-document', methods=['POST'])
def translate_document_endpoint():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    try:
        content = extract_text_from_pdf(file)
        translated_text = translate_text(content)
        return jsonify({'translated_text': translated_text})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Endpoint for image translation
@app.route('/translate-image', methods=['POST'])
def translate_image_endpoint():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400

    image_file = request.files['file']
    if image_file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    try:
        extracted_text = extract_text_from_image(image_file)
        translated_text = translate_text(extracted_text)
        return jsonify({'translated_text': translated_text})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Endpoint for PowerPoint translation
@app.route('/translate-ppt', methods=['POST'])
def translate_ppt_endpoint():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400

    ppt_file = request.files['file']
    if ppt_file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    try:
        # Translate the PPT and return it as a file
        translated_ppt = translate_ppt(ppt_file, from_language="en", to_language="mr")
        return send_file(translated_ppt, as_attachment=True, download_name="translated_ppt.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)
